from datetime import datetime
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt

from models.presentation_content import (
    PresentationContent,
)


class PresentationExportService:
    """Exports generated presentations to PowerPoint."""

    EXPORT_DIR = Path("data/presentations")

    def __init__(self):
        self.EXPORT_DIR.mkdir(
            parents=True,
            exist_ok=True,
        )

    def export(
        self,
        presentation_content: PresentationContent,
        document_name: str,
    ) -> Path:
        """Export presentation content as a PowerPoint file."""

        if not presentation_content.slides:
            raise ValueError(
                "Cannot export an empty presentation."
            )

        if not document_name:
            raise ValueError(
                "Document name is required."
            )

        presentation = Presentation()

        # Widescreen 16:9 layout
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        # Create title slide
        self._add_title_slide(
            presentation=presentation,
            title=presentation_content.title,
            document_name=document_name,
        )

        # Create content slides
        total_slides = len(
            presentation_content.slides
        )

        for position, slide_content in enumerate(
            presentation_content.slides,
            start=1,
        ):

            self._add_content_slide(
                presentation=presentation,
                title=slide_content.title,
                content=slide_content.content,
                position=position,
                total_slides=total_slides,
            )

        # Generate filename
        export_path = self._build_export_path(
            document_name=document_name,
        )

        presentation.save(
            str(export_path)
        )

        return export_path

    def _add_title_slide(
        self,
        presentation: Presentation,
        title: str,
        document_name: str,
    ) -> None:
        """Add the presentation title slide."""

        slide_layout = presentation.slide_layouts[0]

        slide = presentation.slides.add_slide(
            slide_layout
        )

        title_placeholder = slide.shapes.title

        title_placeholder.text = title

        title_paragraph = (
            title_placeholder.text_frame.paragraphs[0]
        )

        title_paragraph.font.size = Pt(32)

        subtitle_placeholder = (
            slide.placeholders[1]
        )

        subtitle_placeholder.text = (
            f"Generated from {document_name}"
        )

        subtitle_paragraph = (
            subtitle_placeholder
            .text_frame
            .paragraphs[0]
        )

        subtitle_paragraph.font.size = Pt(18)

    def _add_content_slide(
        self,
        presentation: Presentation,
        title: str,
        content: list[str],
        position: int,
        total_slides: int,
    ) -> None:
        """Add one content slide."""

        slide_layout = presentation.slide_layouts[1]

        slide = presentation.slides.add_slide(
            slide_layout
        )

        # Slide title
        title_placeholder = slide.shapes.title

        title_placeholder.text = title

        title_paragraph = (
            title_placeholder.text_frame.paragraphs[0]
        )

        title_paragraph.font.size = Pt(28)

        # Slide content
        content_placeholder = slide.placeholders[1]

        text_frame = (
            content_placeholder.text_frame
        )

        text_frame.clear()

        for index, bullet in enumerate(content):

            if index == 0:

                paragraph = (
                    text_frame.paragraphs[0]
                )

            else:

                paragraph = (
                    text_frame.add_paragraph()
                )

            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)

        # Slide number
        slide_number = (
            f"{position} / {total_slides}"
        )

        text_box = slide.shapes.add_textbox(
            Inches(11.8),
            Inches(7.0),
            Inches(1.0),
            Inches(0.3),
        )

        paragraph = (
            text_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = slide_number
        paragraph.font.size = Pt(10)

    def _build_export_path(
        self,
        document_name: str,
    ) -> Path:
        """Build a unique export file path."""

        clean_name = Path(
            document_name
        ).stem

        clean_name = re.sub(
            r"[^a-zA-Z0-9_-]",
            "_",
            clean_name,
        )

        timestamp = datetime.now().strftime(
            "%Y%m%d_%H%M%S"
        )

        filename = (
            f"{clean_name}_presentation_"
            f"{timestamp}.pptx"
        )

        return self.EXPORT_DIR / filename